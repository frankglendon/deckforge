# -*- coding: utf-8 -*-
"""
Vector frameworks: KPI cards / 2x2 maps and chevron value-chains.

These are drawn with native shapes (fully editable) and themed from the palette.
Use them in the left visual area of a content page, or full-width.
"""
import io
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from PIL import Image, ImageOps
from .text import set_run_font

EMU_IN = 914400


def _IN(v):
    return int(v * EMU_IN)


def _fit_crop(path, w_in, h_in):
    """Center-crop an image to the cell aspect (fill, no distortion) in-memory;
    returns a BytesIO PNG/JPEG stream. Never mutates the source file."""
    img = Image.open(path).convert("RGB")
    img = ImageOps.fit(img, (max(1, int(w_in * 200)), max(1, int(h_in * 200))),
                       Image.Resampling.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, "JPEG", quality=90)
    buf.seek(0)
    return buf


def kpi_grid(slide, theme, x, y, w, h, cards, *, cols=2, gap=0.18):
    """cards: list of (big, label, caption). Colored cards with a big figure."""
    n = len(cards)
    rows = (n + cols - 1) // cols
    cw = (w - gap * (cols - 1)) / cols
    ch = (h - gap * (rows - 1)) / rows
    for i, (big, label, caption) in enumerate(cards):
        cx = x + (i % cols) * (cw + gap)
        cy = y + (i // cols) * (ch + gap)
        color = theme.secondary[i % len(theme.secondary)]
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _IN(cx), _IN(cy), _IN(cw), _IN(ch))
        card.fill.solid()
        card.fill.fore_color.rgb = theme.white
        card.line.color.rgb = theme.light_grey
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _IN(cx), _IN(cy), _IN(cw), _IN(0.07))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        bar.shadow.inherit = False
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = _IN(0.16)
        tf.margin_right = _IN(0.14)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        rb = p.add_run()
        rb.text = str(big)
        set_run_font(rb, 24, color, theme=theme, bold=True, serif=True)
        p2 = tf.add_paragraph()
        rl = p2.add_run()
        rl.text = label
        set_run_font(rl, 12, theme.title, theme=theme, bold=True)
        if caption:
            p3 = tf.add_paragraph()
            rc = p3.add_run()
            rc.text = caption
            set_run_font(rc, 9.5, theme.grey, theme=theme)


def chevron(slide, theme, x, y, w, h, steps, *, captions=None):
    """Horizontal value-chain of chevron arrows. steps: list of short labels.
    captions: optional list of one-line descriptions drawn under each step."""
    n = len(steps)
    gap = 0.04
    cw = (w - gap * (n - 1)) / n
    for i, label in enumerate(steps):
        cx = x + i * (cw + gap)
        color = theme.secondary[i % len(theme.secondary)]
        sp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, _IN(cx), _IN(y), _IN(cw), _IN(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
        sp.line.fill.background()
        sp.shadow.inherit = False
        tf = sp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_run_font(r, 11.5, theme.white, theme=theme, bold=True)
        if captions and i < len(captions) and captions[i]:
            cb = slide.shapes.add_textbox(_IN(cx), _IN(y + h + 0.08), _IN(cw), _IN(0.9))
            cb.text_frame.word_wrap = True
            cp = cb.text_frame.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            rc = cp.add_run()
            rc.text = captions[i]
            set_run_font(rc, 9.5, theme.text, theme=theme)


def image_row(slide, paths, x, y, w, *, gap=0.12, ratio=2.4, captions=None):
    """Place a row of pre-cropped thumbnails (each `ratio` wide:tall). Returns
    the row height used. Thumbnails should already be cropped to `ratio`."""
    n = len(paths)
    if n == 0:
        return 0
    cw = (w - gap * (n - 1)) / n
    ch = cw / ratio
    for i, p in enumerate(paths):
        pic = _fit_crop(p, cw, ch)  # center-crop to the exact cell aspect (no distortion)
        slide.shapes.add_picture(pic, _IN(x + i * (cw + gap)), _IN(y), _IN(cw), _IN(ch))
        if captions and i < len(captions) and captions[i]:
            cb = slide.shapes.add_textbox(_IN(x + i * (cw + gap)), _IN(y + ch + 0.02),
                                          _IN(cw), _IN(0.3))
            cp = cb.text_frame.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            rc = cp.add_run()
            rc.text = captions[i]
            set_run_font(rc, 8, theme.grey, theme=theme)
    return ch
