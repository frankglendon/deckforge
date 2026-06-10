# -*- coding: utf-8 -*-
"""
Native PowerPoint charts (column / bar / line / doughnut).

Everything is a real, editable PowerPoint chart object (NOT a matplotlib image)
— so anyone opening the deck can edit the data, colors and labels. One point
per chart may be highlighted in the accent color (the single key datapoint).
"""
from pptx.util import Emu, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

EMU_IN = 914400


def _IN(v):
    return int(v * EMU_IN)


def _style(chart, theme, legend=False,
           legend_pos=XL_LEGEND_POSITION.BOTTOM):
    try:
        chart.has_title = False
    except Exception:
        pass
    chart.has_legend = legend
    if legend:
        chart.legend.position = legend_pos
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = theme.font_body
    for ax in ("category_axis", "value_axis"):
        try:
            a = getattr(chart, ax)
            a.tick_labels.font.size = Pt(theme.size_chart_label)
            a.tick_labels.font.name = theme.font_body
        except Exception:
            pass


def _label_font(plot, theme):
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(theme.size_chart_label)
    dl.font.bold = True
    dl.font.name = theme.font_body
    dl.font.color.rgb = theme.text
    return dl


def column(slide, theme, x, y, w, h, categories, series, *,
           highlight=None, number_format=None, legend=None):
    """series: dict {name: [values]}. highlight: (series_idx, point_idx) -> accent."""
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series.items():
        cd.add_series(name, vals, number_format=number_format)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                   _IN(x), _IN(y), _IN(w), _IN(h), cd).chart
    plot = chart.plots[0]
    for si, ser in enumerate(plot.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = theme.secondary[si % len(theme.secondary)]
    if highlight is not None:
        si, pi = highlight
        pt = plot.series[si].points[pi]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = theme.accent
    dl = _label_font(plot, theme)
    try:
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
    except Exception:
        pass
    _style(chart, theme, legend=legend if legend is not None else len(series) > 1)
    return chart


def bar(slide, theme, x, y, w, h, categories, values, *, highlight_idx=None):
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                                   _IN(x), _IN(y), _IN(w), _IN(h), cd).chart
    ser = chart.plots[0].series[0]
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = theme.secondary[0]
    if highlight_idx is not None:
        pt = ser.points[highlight_idx]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = theme.accent
    _label_font(chart.plots[0], theme)
    _style(chart, theme, legend=False)
    return chart


def line(slide, theme, x, y, w, h, categories, series):
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series.items():
        cd.add_series(name, vals)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS,
                                   _IN(x), _IN(y), _IN(w), _IN(h), cd).chart
    for si, s in enumerate(chart.plots[0].series):
        s.format.line.color.rgb = theme.secondary[si % len(theme.secondary)]
        s.format.line.width = Pt(2.5)
        s.smooth = False
    _style(chart, theme, legend=len(series) > 1)
    return chart


def doughnut(slide, theme, x, y, w, h, categories, values, *, highlight_idx=None):
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("", values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT,
                                   _IN(x), _IN(y), _IN(w), _IN(h), cd).chart
    pts = chart.plots[0].series[0].points
    for i in range(len(values)):
        c = theme.accent if i == highlight_idx else theme.secondary[i % len(theme.secondary)]
        pts[i].format.fill.solid()
        pts[i].format.fill.fore_color.rgb = c
    dl = chart.plots[0].data_labels
    chart.plots[0].has_data_labels = True
    dl.number_format = "0%"
    dl.number_format_is_linked = False
    dl.show_percentage = True
    dl.show_value = False
    dl.font.size = Pt(theme.size_chart_label)
    dl.font.bold = True
    dl.font.name = theme.font_body
    # NOTE: do NOT set dl.position for pie/doughnut — PowerPoint flags it for repair.
    _style(chart, theme, legend=True, legend_pos=XL_LEGEND_POSITION.RIGHT)
    return chart
