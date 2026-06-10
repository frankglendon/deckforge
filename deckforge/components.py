# -*- coding: utf-8 -*-
"""
Insight blocks — the atomic right-hand column of a content slide.

Three stacked blocks mirror the consulting spine: Facts -> Insights ->
Implications. Each block = tinted background + left accent bar + a bold tag
(16pt) + rich body (12pt, 1.5 line spacing, key terms/numbers auto-highlighted).
"""
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from .text import set_run_font, add_rich

EMU_IN = 914400


def _IN(v):
    return int(v * EMU_IN)


def insight_block(slide, theme, x, y, w, h, tag, body, accent, tint):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _IN(x), _IN(y), _IN(w), _IN(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = tint
    bg.line.fill.background()
    bg.shadow.inherit = False
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _IN(x), _IN(y), _IN(0.06), _IN(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bg.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = _IN(0.16)
    tf.margin_right = _IN(0.14)
    tf.margin_top = _IN(0.09)
    tf.margin_bottom = _IN(0.07)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(3)
    rt = p.add_run()
    rt.text = tag
    set_run_font(rt, theme.size_block_tag, accent, theme=theme, bold=True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.line_spacing = 1.5
    add_rich(p2, body, theme.size_body, theme=theme, base=theme.text)
    return bg


def insight_col(slide, theme, x, y, w, h, blocks, gap=0.14):
    """blocks: list of (tag, body). Up to 3 typical (Facts/Insights/Implications)."""
    n = len(blocks)
    bh = (h - gap * (n - 1)) / n
    palette = theme.block_themes()
    for i, (tag, body) in enumerate(blocks):
        accent, tint = palette[i % len(palette)]
        insight_block(slide, theme, x, y + i * (bh + gap), w, bh, tag, body, accent, tint)
