# -*- coding: utf-8 -*-
"""
Text helpers: font application (serif headings / sans body, each with an
East-Asian typeface + language tag) and "rich emphasis" — auto-highlighting
key terms and numbers in the accent color.

Why the explicit <a:ea> / <a:latin> dance: python-pptx's run.font.name only sets
the Latin typeface. Without an explicit East-Asian (`ea`) typeface, CJK glyphs
fall back to a default font and render as tofu boxes on some platforms. We set
both, plus a `lang` attribute so spellcheck/shaping behave.
"""
import re
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# Match quoted key terms 「...」/"..." and numbers with a unit -> emphasize.
EMPH = re.compile(
    r'(「[^」]*」|“[^”]*”|[+\-]?[0-9][0-9.,/]*(?:%|x|×|bn|m|k|B|M|K|pp|pts?|'
    r'亿|万|倍|个|年|岁)?)'
)


def set_run_font(run, size, color, *, theme, bold=False, underline=False,
                 italic=False, serif=False):
    """Apply size/color/bold + Latin & East-Asian typefaces + language tag.
    serif=True uses the heading (serif) font pair; otherwise the body (sans)."""
    latin, ea = theme.fonts(serif=serif)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline
    run.font.color.rgb = color
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", theme.lang)
    for tag in ("a:ea", "a:latin"):
        for old in rPr.findall(qn(tag)):
            rPr.remove(old)
    rPr.append(rPr.makeelement(qn("a:latin"), {"typeface": latin}))
    rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": ea}))


def add_rich(paragraph, text, size, *, theme, base=None, emphasize=True, serif=False):
    """Split text on the emphasis regex and color key terms/numbers in accent."""
    base = base if base is not None else theme.text
    if not emphasize or not getattr(theme, "accent_text", False):
        r = paragraph.add_run()
        r.text = text
        set_run_font(r, size, base, theme=theme, serif=serif)
        return
    for seg in EMPH.split(text):
        if not seg:
            continue
        is_key = (bool(EMPH.fullmatch(seg)) and any(c.isdigit() for c in seg)) or \
            seg.startswith(("「", "“"))
        r = paragraph.add_run()
        r.text = seg
        set_run_font(r, size, theme.primary if is_key else base,
                     theme=theme, bold=is_key, serif=serif)


def set_text(text_frame, text, *, theme, size=12, color=None, bold=False,
             align=PP_ALIGN.LEFT, serif=False):
    color = color if color is not None else theme.text
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run_font(r, size, color, theme=theme, bold=bold, serif=serif)
    return text_frame
