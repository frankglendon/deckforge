# -*- coding: utf-8 -*-
"""
Comparison tables — the most persuasive way to contrast competitors/options.

Themed header row, zebra body, thin borders (written in schema order so
PowerPoint doesn't flag the file for repair), and an optional bottom
"takeaway / implication" bar.
"""
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from .text import set_run_font, add_rich

EMU_IN = 914400


def _IN(v):
    return int(v * EMU_IN)


def _cell_border(cell, theme, w=9525):
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
    hexv = "%02X%02X%02X" % (theme.light_grey[0], theme.light_grey[1], theme.light_grey[2])
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):  # schema order matters
        ln = tcPr.makeelement(qn(tag), {"w": str(w), "cap": "flat"})
        fill = ln.makeelement(qn("a:solidFill"), {})
        clr = fill.makeelement(qn("a:srgbClr"), {"val": hexv})
        fill.append(clr)
        ln.append(fill)
        tcPr.append(ln)


def table(slide, theme, x, y, w, h, rows, *, col_widths=None, zebra=True):
    """rows: 2D list incl. header row."""
    nrow, ncol = len(rows), len(rows[0])
    gf = slide.shapes.add_table(nrow, ncol, _IN(x), _IN(y), _IN(w), _IN(h))
    tbl = gf.table
    tbl.first_row = False  # disable template styling; we control colors
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = _IN(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = _IN(0.06)
            cell.margin_right = _IN(0.06)
            cell.margin_top = _IN(0.03)
            cell.margin_bottom = _IN(0.03)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if (ci > 0 or ri == 0) else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = theme.primary
                set_run_font(r, theme.size_table_header, theme.white, theme=theme, bold=True)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = theme.light_grey if (zebra and ri % 2 == 0) else theme.white
                set_run_font(r, theme.size_table_body, theme.text, theme=theme)
            _cell_border(cell, theme)
    return tbl


def takeaway_bar(slide, theme, x, y, w, text, label="Implication"):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _IN(x), _IN(y), _IN(w), _IN(0.55))
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme.tint_primary
    bg.line.fill.background()
    bg.shadow.inherit = False
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _IN(x), _IN(y), _IN(0.06), _IN(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.primary
    bar.line.fill.background()
    bar.shadow.inherit = False
    tb = slide.shapes.add_textbox(_IN(x + 0.2), _IN(y), _IN(w - 0.3), _IN(0.55))
    tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tb.text_frame.paragraphs[0]
    lead = p.add_run()
    lead.text = f"{label}:  "
    set_run_font(lead, 12, theme.primary, theme=theme, bold=True)
    add_rich(p, text, 12, theme=theme, base=theme.text)
    return bg
