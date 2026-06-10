# -*- coding: utf-8 -*-
"""
Deck: the high-level builder.

Layout faithfully replicates a McKinsey-style executive report (calibrated to a
public "Global Economics Intelligence" summary): full-width serif headline in
black with a hairline rule beneath, a small grey kicker top-right, a left visual
+ right justified body column, and a bottom hairline with "<brand>  <page>".

No proprietary template, fonts or trademarks are bundled — fonts are generic
substitutes (Georgia / KaiTi / Arial) and the brand string is yours to set.

Builds on a BLANK presentation by drawing every element at absolute positions,
so it is fully self-contained and reproducible. Swap `Theme` to re-brand.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR

from .theme import Theme, DEFAULT
from .text import set_run_font, add_rich, EMPH
from .components import insight_col

EMU_IN = 914400


def IN(v):
    return int(v * EMU_IN)


DEFAULT_LABELS = {
    "source": "Source: ",
    "part": "PART {n}",
    "in_this_part": "In this part",
    "implication": "Implication",
    "cover_note": "Design language inspired by public consulting reports. "
                  "Sample / illustrative.",
}


class Deck:
    def __init__(self, theme: Theme = DEFAULT, brand="Acme Research",
                 footer="Acme Research", labels=None):
        self.theme = theme
        self.brand = brand
        self.footer_text = footer
        self.labels = dict(DEFAULT_LABELS, **(labels or {}))  # i18n UI strings
        self.prs = Presentation()
        self.prs.slide_width = IN(13.333)
        self.prs.slide_height = IN(7.5)
        self._blank = self.prs.slide_layouts[6]  # "Blank" in the default template
        self.page = 0
        self.image_pool = []   # thumbnails (cropped ~2.4:1) auto-placed on content pages
        self._pi = 0

    def set_image_pool(self, paths):
        """Provide a list of pre-cropped (~2.4:1) thumbnail paths; content pages
        will draw a 2-photo row from this pool, cycling with spacing."""
        self.image_pool = list(paths)
        self._pi = 0

    def _take_images(self, k):
        if not self.image_pool:
            return []
        pool = self.image_pool
        out = [pool[(self._pi + j) % len(pool)] for j in range(min(k, len(pool)))]
        self._pi = (self._pi + k) % len(pool)
        return out

    # ---------- low-level primitives ----------
    def _slide(self):
        self.page += 1
        return self.prs.slides.add_slide(self._blank)

    def rect(self, slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
        sp = slide.shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            sp.line.width = Pt(0.75)
        sp.shadow.inherit = False
        return sp

    def hrule(self, slide, x, y, w, color=None, weight=0.75):
        cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x), IN(y),
                                        IN(x + w), IN(y))
        cn.line.color.rgb = color or self.theme.rule
        cn.line.width = Pt(weight)
        cn.shadow.inherit = False
        return cn

    def textbox(self, slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        return tb

    def _para(self, tf, first=False):
        return tf.paragraphs[0] if first else tf.add_paragraph()

    # ---------- shared furniture (GEI signature) ----------
    def kicker(self, slide, text):
        """Small grey tag, top-right (e.g. 'Section | 15 April 2026')."""
        tb = self.textbox(slide, 6.5, 0.26, 6.4, 0.3)
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = text
        set_run_font(r, self.theme.size_footnote, self.theme.grey, theme=self.theme)
        return tb

    def action_title(self, slide, text, y=0.46, rule=True):
        """Full-sentence black serif headline (a complete cause -> effect / so-what
        statement, typically two lines) + hairline rule beneath (GEI style)."""
        tb = self.textbox(slide, 0.5, y, 12.3, 1.05)
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if not self.theme.accent_text:
            # McKinsey-faithful: monochrome black serif headline
            r = p.add_run()
            r.text = text
            set_run_font(r, self.theme.size_title, self.theme.title,
                         theme=self.theme, bold=True, serif=True)
        else:
            for seg in EMPH.split(text):
                if not seg:
                    continue
                is_key = bool(EMPH.fullmatch(seg)) and (
                    any(c.isdigit() for c in seg) or seg.startswith(("「", "“")))
                r = p.add_run()
                r.text = seg
                set_run_font(r, self.theme.size_title,
                             self.theme.primary if is_key else self.theme.title,
                             theme=self.theme, bold=True, serif=True)
        if rule:
            self.hrule(slide, 0.5, 1.5, 12.33)
        return tb

    def subhead(self, slide, text, y=1.6):
        tb = self.textbox(slide, 0.5, y, 12.3, 0.45)
        add_rich(tb.text_frame.paragraphs[0], text, 11,
                 theme=self.theme, base=self.theme.grey, emphasize=False)
        return tb

    def source_line(self, slide, sources, y=6.92):
        if not sources:
            return
        tb = self.textbox(slide, 0.5, y, 9.5, 0.3)
        p = tb.text_frame.paragraphs[0]
        lead = p.add_run()
        lead.text = self.labels["source"]
        set_run_font(lead, self.theme.size_footnote, self.theme.grey, theme=self.theme)
        for i, (label, url) in enumerate(sources):
            if i:
                sep = p.add_run()
                sep.text = "; "
                set_run_font(sep, self.theme.size_footnote, self.theme.grey, theme=self.theme)
            r = p.add_run()
            r.text = label
            if url:
                r.hyperlink.address = url
                set_run_font(r, self.theme.size_footnote, self.theme.primary,
                             theme=self.theme, underline=True)
            else:
                set_run_font(r, self.theme.size_footnote, self.theme.grey, theme=self.theme)

    def footer(self, slide):
        self.hrule(slide, 0.5, 7.18, 12.33, color=self.theme.rule, weight=0.5)
        tb = self.textbox(slide, 9.0, 7.2, 3.83, 0.28)
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = f"{self.footer_text}    {self.page}"
        set_run_font(r, 9, self.theme.text, theme=self.theme)

    def _finish(self, slide, sources):
        self.source_line(slide, sources)
        self.footer(slide)

    # ---------- right-hand body column (GEI flowing text) ----------
    def _body_column(self, slide, x, y, w, h, paragraphs):
        """paragraphs: list of str OR (bold_lead, rest). Justified Arial body."""
        tb = self.textbox(slide, x, y, w, h)
        tf = tb.text_frame
        for i, para in enumerate(paragraphs):
            p = self._para(tf, first=(i == 0))
            p.alignment = PP_ALIGN.JUSTIFY
            p.line_spacing = 1.12
            p.space_after = Pt(8)
            if isinstance(para, (tuple, list)):
                lead, rest = para
                rl = p.add_run()
                rl.text = lead + " "
                set_run_font(rl, self.theme.size_body, self.theme.text,
                             theme=self.theme, bold=True)
                add_rich(p, rest, self.theme.size_body, theme=self.theme)
            else:
                add_rich(p, para, self.theme.size_body, theme=self.theme)
        return tb

    def chart_caption(self, slide, x, y, w, title, unit):
        """GEI chart caption: bold title + light unit line above the chart."""
        tb = self.textbox(slide, x, y, w, 0.5)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_run_font(r, 11, self.theme.title, theme=self.theme, bold=True)
        if unit:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = unit
            set_run_font(r2, 9, self.theme.grey, theme=self.theme, italic=True)
        return tb

    # ---------- page types ----------
    def cover(self, title, subtitle="", meta="", image_path=None, cta=None):
        slide = self._slide()
        self.rect(slide, 0, 0, 13.333, 7.5, self.theme.dark_bg)
        if image_path:
            slide.shapes.add_picture(image_path, IN(8.0), 0, IN(5.33), IN(7.5))
        else:
            # signature thin-line motif fills the right half (deep blue + cyan)
            self.line_pattern(slide, 8.7, 2.2, 3.9, n=22, dy=0.12, weight=1.4)
        bt = self.textbox(slide, 0.6, 0.45, 6, 0.6)
        for i, line in enumerate(self.brand.split("\n")):
            p = self._para(bt.text_frame, first=(i == 0))
            r = p.add_run()
            r.text = line
            set_run_font(r, 14, self.theme.white, theme=self.theme, bold=True)
        tt = self.textbox(slide, 0.6, 2.5, 7.0, 2.2)
        for i, line in enumerate(title.split("\n")):
            p = self._para(tt.text_frame, first=(i == 0))
            run = p.add_run()
            run.text = line
            set_run_font(run, 36, self.theme.white, theme=self.theme, bold=True, serif=True)
        if subtitle:
            st = self.textbox(slide, 0.62, 4.55, 7.0, 0.6)
            r = st.text_frame.paragraphs[0].add_run()
            r.text = subtitle
            set_run_font(r, 16, self.theme.light_grey, theme=self.theme)
        if meta:
            mt = self.textbox(slide, 0.62, 5.15, 7.0, 0.5)
            r = mt.text_frame.paragraphs[0].add_run()
            r.text = meta
            set_run_font(r, 11, self.theme.light_grey, theme=self.theme)
        if cta:
            self.rect(slide, 0.62, 5.95, 1.7, 0.45, self.theme.primary)
            cb = self.textbox(slide, 0.62, 5.97, 1.7, 0.42, anchor=MSO_ANCHOR.MIDDLE)
            cp = cb.text_frame.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            r = cp.add_run()
            r.text = cta
            set_run_font(r, 11, self.theme.white, theme=self.theme, bold=True)
        dt = self.textbox(slide, 0.62, 6.95, 7.5, 0.4)
        r = dt.text_frame.paragraphs[0].add_run()
        r.text = self.labels["cover_note"]
        set_run_font(r, 7, self.theme.grey, theme=self.theme)
        return slide

    def content(self, title, body, *, kicker=None, subhead=None, sources=None,
                images=None):
        """McKinsey-style content page: title + rule, LEFT visual (rect returned),
        RIGHT justified body text. A row of related photos is added below the left
        visual (from the deck image pool, or pass images=[...]; images=False off).
        Returns (slide, left_rect_inches)."""
        slide = self._slide()
        if kicker:
            self.kicker(slide, kicker)
        self.action_title(slide, title)
        if subhead:
            self.subhead(slide, subhead)
        self._body_column(slide, 7.0, 1.75, 5.85, 5.0, body)
        left = self._left_with_images(slide, images)
        self._finish(slide, sources or [])
        return slide, left

    def content_blocks(self, title, blocks, *, kicker=None, subhead=None,
                       sources=None, images=None):
        """Variant: title + rule, LEFT visual + RIGHT three insight blocks
        (Facts / Insights / Implications). Returns (slide, left_rect)."""
        slide = self._slide()
        if kicker:
            self.kicker(slide, kicker)
        self.action_title(slide, title)
        if subhead:
            self.subhead(slide, subhead)
        insight_col(slide, self.theme, x=7.0, y=1.8, w=5.85, h=4.9, blocks=blocks)
        left = self._left_with_images(slide, images)
        self._finish(slide, sources or [])
        return slide, left

    def _left_with_images(self, slide, images):
        """Place a 3-photo row low in the left column; return the (shorter)
        left visual rect above it. If no images, return the full-height rect."""
        imgs = images if images is not None else self._take_images(3)
        if imgs:
            from .frameworks import image_row
            image_row(slide, imgs, 0.5, 5.18, 6.1, gap=0.1)
            return (0.5, 1.78, 6.1, 3.1)
        return (0.5, 1.8, 6.1, 4.9)

    def wide_slide(self, title, *, kicker=None, subhead=None, sources=None,
                   images=None):
        """Full-width content page (tables / big frameworks). A related-photo row
        is added at the bottom (from the pool, or pass images=[...] / images=False).
        Returns (slide, content_rect)."""
        slide = self._slide()
        if kicker:
            self.kicker(slide, kicker)
        self.action_title(slide, title)
        if subhead:
            self.subhead(slide, subhead)
        imgs = images if images is not None else self._take_images(4)
        if imgs:
            from .frameworks import image_row
            image_row(slide, imgs, 0.5, 5.52, 12.33)
            rect = (0.5, 1.85, 12.33, 3.05)
        else:
            rect = (0.5, 1.85, 12.33, 4.7)
        self._finish(slide, sources or [])
        return slide, rect

    def core_questions(self, headline, intro, cards):
        """cards: list of (number, question, subtitle, [bullets])."""
        slide = self._slide()
        self.action_title(slide, headline)
        if intro:
            self.subhead(slide, intro)
        n = len(cards)
        gap = 0.3
        cw = (12.33 - gap * (n - 1)) / n
        x = 0.5
        for (num, q, sub, bullets) in cards:
            self._question_card(slide, x, 2.0, cw, 4.7, num, q, sub, bullets)
            x += cw + gap
        self.footer(slide)
        return slide

    def _question_card(self, slide, x, y, w, h, num, q, sub, bullets):
        self.rect(slide, x, y, w, h, self.theme.white, line=self.theme.light_grey)
        self.rect(slide, x, y, w, 0.05, self.theme.primary)
        nb = self.textbox(slide, x + 0.18, y + 0.16, 1, 0.4)
        r = nb.text_frame.paragraphs[0].add_run()
        r.text = str(num)
        set_run_font(r, 18, self.theme.grey, theme=self.theme, bold=True, serif=True)
        qb = self.textbox(slide, x + 0.18, y + 0.62, w - 0.36, 0.7)
        rq = qb.text_frame.paragraphs[0].add_run()
        rq.text = q
        set_run_font(rq, 15, self.theme.title, theme=self.theme, bold=True, serif=True)
        if sub:
            sb = self.textbox(slide, x + 0.2, y + 1.4, w - 0.4, 0.4)
            rs = sb.text_frame.paragraphs[0].add_run()
            rs.text = sub
            set_run_font(rs, 11, self.theme.primary, theme=self.theme, bold=True)
        bt = self.textbox(slide, x + 0.2, y + 1.9, w - 0.4, h - 2.05)
        for i, b in enumerate(bullets):
            p = self._para(bt.text_frame, first=(i == 0))
            p.space_after = Pt(6)
            bullet = p.add_run()
            bullet.text = "— "
            set_run_font(bullet, 12, self.theme.primary, theme=self.theme, bold=True)
            add_rich(p, b, 12, theme=self.theme)

    def agenda(self, headline, items):
        """items: list of (number, chapter_one_liner)."""
        slide = self._slide()
        self.action_title(slide, headline)
        y = 2.0
        rowh = (6.7 - y) / len(items)
        for i, (num, line) in enumerate(items):
            yy = y + i * rowh
            self.rect(slide, 0.5, yy, 0.75, rowh - 0.2,
                      self.theme.secondary[i % len(self.theme.secondary)])
            nb = self.textbox(slide, 0.5, yy, 0.75, rowh - 0.2, anchor=MSO_ANCHOR.MIDDLE)
            np_ = nb.text_frame.paragraphs[0]
            np_.alignment = PP_ALIGN.CENTER
            rn = np_.add_run()
            rn.text = f"{num:02d}"
            set_run_font(rn, 16, self.theme.white, theme=self.theme, bold=True, serif=True)
            tb = self.textbox(slide, 1.5, yy, 11.3, rowh - 0.2, anchor=MSO_ANCHOR.MIDDLE)
            add_rich(tb.text_frame.paragraphs[0], line, 14, theme=self.theme,
                     base=self.theme.title)
        self.footer(slide)
        return slide

    def section_divider(self, part_no, title, subtitle, toc_items,
                        takeaway="", image_path=None):
        slide = self._slide()
        self.rect(slide, 0, 0, 13.333, 7.5, self.theme.dark_bg)
        if image_path:
            slide.shapes.add_picture(image_path, IN(8.4), IN(1.2), IN(4.4), IN(3.4))
        pt = self.textbox(slide, 0.6, 0.6, 7, 0.4)
        rp = pt.text_frame.paragraphs[0].add_run()
        rp.text = self.labels["part"].format(n=part_no)
        set_run_font(rp, 13, self.theme.primary, theme=self.theme, bold=True)
        tt = self.textbox(slide, 0.6, 1.1, 7.5, 1.2)
        rt = tt.text_frame.paragraphs[0].add_run()
        rt.text = title
        set_run_font(rt, 32, self.theme.white, theme=self.theme, bold=True, serif=True)
        if subtitle:
            st = self.textbox(slide, 0.6, 2.15, 7.5, 0.45)
            rs = st.text_frame.paragraphs[0].add_run()
            rs.text = subtitle
            set_run_font(rs, 13, self.theme.light_grey, theme=self.theme)
        hb = self.textbox(slide, 0.6, 2.8, 7.6, 0.35)
        rh = hb.text_frame.paragraphs[0].add_run()
        rh.text = self.labels["in_this_part"]
        set_run_font(rh, 12, self.theme.primary, theme=self.theme, bold=True)
        self.hrule(slide, 0.6, 3.16, 1.4, color=self.theme.primary, weight=1.5)
        lb = self.textbox(slide, 0.6, 3.3, 7.7, 3.0)
        tf = lb.text_frame
        for i, it in enumerate(toc_items):
            p = self._para(tf, first=(i == 0))
            p.line_spacing = 1.5
            mark = p.add_run()
            mark.text = "—  "
            set_run_font(mark, self.theme.size_toc, self.theme.primary, theme=self.theme, bold=True)
            body = p.add_run()
            body.text = it
            set_run_font(body, self.theme.size_toc, self.theme.white, theme=self.theme)
        if takeaway:
            self.rect(slide, 0.6, 6.5, 12.13, 0.6, RGBColor(0x0B, 0x2A, 0x3A))
            self.rect(slide, 0.6, 6.5, 0.06, 0.6, self.theme.primary)
            kb = self.textbox(slide, 0.8, 6.55, 11.8, 0.5, anchor=MSO_ANCHOR.MIDDLE)
            kp = kb.text_frame.paragraphs[0]
            lead = kp.add_run()
            lead.text = self.labels["implication"] + "   "
            set_run_font(lead, 12, self.theme.primary, theme=self.theme, bold=True)
            add_rich(kp, takeaway, 12, theme=self.theme, base=self.theme.light_grey)
        return slide

    def prose(self, title, paragraphs, *, kicker=None, sources=None, columns=1,
              image_path=None):
        """Full-width text page (disclaimer / exec summary / section summary /
        methodology). paragraphs: list of str OR (bold_lead, rest).
        image_path: optional photo placed on the right (forces single column)."""
        slide = self._slide()
        if kicker:
            self.kicker(slide, kicker)
        self.action_title(slide, title)
        if image_path:
            self._body_column(slide, 0.5, 1.85, 6.7, 4.9, paragraphs)
            slide.shapes.add_picture(image_path, IN(7.55), IN(1.95), IN(5.28), IN(4.32))
        elif columns == 2:
            half = len(paragraphs) // 2 + len(paragraphs) % 2
            self._body_column(slide, 0.5, 1.85, 6.0, 4.9, paragraphs[:half])
            self._body_column(slide, 6.83, 1.85, 6.0, 4.9, paragraphs[half:])
        else:
            self._body_column(slide, 0.5, 1.85, 12.33, 4.9, paragraphs)
        self._finish(slide, sources or [])
        return slide

    def source_list(self, title, sources_by_group, *, kicker=None):
        """Methodology/sources page: groups of hyperlinked sources in two columns."""
        slide = self._slide()
        if kicker:
            self.kicker(slide, kicker)
        self.action_title(slide, title)
        groups = list(sources_by_group.items())
        mid = len(groups) // 2 + len(groups) % 2
        for col, chunk in enumerate((groups[:mid], groups[mid:])):
            x = 0.5 + col * 6.4
            tb = self.textbox(slide, x, 1.9, 6.1, 4.9)
            tf = tb.text_frame
            first = True
            for gname, items in chunk:
                gp = self._para(tf, first=first)
                first = False
                gp.space_before = Pt(6)
                rg = gp.add_run()
                rg.text = gname
                set_run_font(rg, 11, self.theme.title, theme=self.theme, bold=True)
                for label, url in items:
                    p = self._para(tf)
                    p.line_spacing = 1.15
                    mark = p.add_run()
                    mark.text = "—  "
                    set_run_font(mark, 9, self.theme.grey, theme=self.theme)
                    r = p.add_run()
                    r.text = label
                    if url:
                        r.hyperlink.address = url
                        set_run_font(r, 9, self.theme.primary, theme=self.theme, underline=True)
                    else:
                        set_run_font(r, 9, self.theme.text, theme=self.theme)
        self.footer(slide)
        return slide

    def back_cover(self, line="Thank you", sub=""):
        slide = self._slide()
        self.rect(slide, 0, 0, 13.333, 7.5, self.theme.dark_bg)
        self.line_pattern(slide, 0.7, 5.6, 4.0, n=14, dy=0.11, weight=1.3)
        bt = self.textbox(slide, 0.6, 0.5, 6, 0.6)
        r = bt.text_frame.paragraphs[0].add_run()
        r.text = self.brand
        set_run_font(r, 14, self.theme.white, theme=self.theme, bold=True)
        tt = self.textbox(slide, 0.6, 2.9, 11, 1.4)
        r = tt.text_frame.paragraphs[0].add_run()
        r.text = line
        set_run_font(r, 30, self.theme.white, theme=self.theme, bold=True, serif=True)
        if sub:
            st = self.textbox(slide, 0.62, 4.15, 11, 0.6)
            r = st.text_frame.paragraphs[0].add_run()
            r.text = sub
            set_run_font(r, 12, self.theme.light_grey, theme=self.theme)
        return slide

    def conclusion(self, headline, body_lines, sources=None):
        slide = self._slide()
        self.rect(slide, 0, 0, 13.333, 7.5, self.theme.dark_bg)
        self.rect(slide, 0.5, 1.0, 0.08, 4.6, self.theme.primary)
        tt = self.textbox(slide, 0.85, 1.1, 11.8, 2.2)
        rt = tt.text_frame.paragraphs[0].add_run()
        rt.text = headline
        set_run_font(rt, 26, self.theme.white, theme=self.theme, bold=True, serif=True)
        bt = self.textbox(slide, 0.87, 3.5, 11.6, 3.0)
        for i, line in enumerate(body_lines):
            p = self._para(bt.text_frame, first=(i == 0))
            p.space_after = Pt(8)
            p.line_spacing = 1.3
            mark = p.add_run()
            mark.text = "—  "
            set_run_font(mark, 14, self.theme.primary, theme=self.theme, bold=True)
            add_rich(p, line, 14, theme=self.theme, base=self.theme.light_grey)
        if sources:
            self.source_line(slide, sources)
        return slide

    def line_pattern(self, slide, x, y, w, n=18, color=None, weight=1.1, dy=0.085):
        """McKinsey-style thin parallel line motif (a stack of cyan hairlines)."""
        color = color or self.theme.primary
        # deterministic varied lengths -> a subtle 'line-chart' silhouette
        pat = [0.35, 0.55, 0.45, 0.7, 0.6, 0.85, 0.75, 1.0, 0.9, 0.8,
               0.95, 0.65, 0.78, 0.5, 0.68, 0.42, 0.58, 0.48]
        for i in range(n):
            lw = w * pat[i % len(pat)]
            cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                            IN(x), IN(y + i * dy),
                                            IN(x + lw), IN(y + i * dy))
            cn.line.color.rgb = color
            cn.line.width = Pt(weight)
            cn.shadow.inherit = False

    def save(self, path):
        self.prs.save(path)
        return path
